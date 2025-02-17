from docx import Document
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO


# PPTX Export
def generate_pptx(presentation_data):
    prs = Presentation()

    for slide_data in presentation_data:
        slide_layout = prs.slide_layouts[1]  # Title & Content Layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_data["title"]

        # Add bullet points for better formatting
        content.text = ""
        for point in slide_data["points"]:
            p = content.text_frame.add_paragraph()
            p.text = f"• {point}"
            p.space_after = Inches(0.2)

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# DOCX Export
def generate_docx(content):
    doc = Document()
    doc.add_paragraph(content)
    
    docx_stream = BytesIO()
    doc.save(docx_stream)
    docx_stream.seek(0)
    return docx_stream

# PDF Export
def generate_pdf(content):
    pdf_buffer = BytesIO()
    c = canvas.Canvas(pdf_buffer, pagesize=letter)
    c.setFont("Helvetica", 10)
    
    text = c.beginText(40, 750)
    text.textLines(content)
    
    c.drawText(text)
    c.showPage()
    c.save()

    pdf_buffer.seek(0)
    return pdf_buffer
