import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from utils.ai_helper import generate_presentation
from io import BytesIO

def generate_pptx(presentation_data):
    prs = Presentation()

    if not isinstance(presentation_data, list):  # Ensure it's a list
        st.error("Invalid data format received from AI. Try again.")
        return None

    for slide_data in presentation_data:
        if not isinstance(slide_data, dict) or "title" not in slide_data or "points" not in slide_data:
            continue  # Skip invalid entries

        slide_layout = prs.slide_layouts[1]  # Title & Content Layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_data["title"]

        content.text = ""
        for point in slide_data["points"]:
            p = content.text_frame.add_paragraph()
            p.text = f"• {point}"
            p.space_after = Inches(0.2)

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

def presentation_generator():
    st.subheader("📊 AI-Powered Presentation Generator")
    st.write("Generate structured PowerPoint slides with proper formatting.")

    subject = st.text_input("📘 Subject", help="Enter the subject (e.g., Artificial Intelligence).")
    topic = st.text_input("📌 Topic", help="Enter the topic (e.g., Neural Networks).")
    num_slides = st.number_input("📈 Number of Slides", min_value=1, max_value=50, value=10)
    additional_notes = st.text_area("Additional Notes",placeholder="enter any relevant points you want to include")


    if st.button("Generate Presentation"):
        with st.spinner("Generating slides..."):
            presentation_data = generate_presentation(subject, topic, num_slides, additional_notes)

            if not isinstance(presentation_data, list) or not all(isinstance(slide, dict) for slide in presentation_data):
                st.error("⚠️ Failed to generate structured slides. Try again.")
            else:
                st.session_state["presentation"] = presentation_data
                st.success("✅ Presentation generated successfully!")

    if "presentation" in st.session_state:
        ppt_data = generate_pptx(st.session_state["presentation"])
        st.download_button("📥 Download Presentation (PPTX)", ppt_data, "presentation.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
